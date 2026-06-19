from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ===== Couleurs Ferrari =====
ROSSO = RGBColor(0xDC, 0x00, 0x00)
DARK = RGBColor(0x0A, 0x0A, 0x0A)
PANEL = RGBColor(0x12, 0x12, 0x12)
WHITE = RGBColor(0xF5, 0xF5, 0xF5)
MUTED = RGBColor(0x8A, 0x8A, 0x8A)
GREEN = RGBColor(0x00, 0xFF, 0x41)
AMBER = RGBColor(0xFF, 0xB8, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
        p.level = 0
        # Bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
    return tf

def add_line(slide, left, top, width, color=ROSSO):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_red_bar(slide):
    add_line(slide, 0.5, 0.55, 2, ROSSO)

def add_footer(slide, text="Scuderia Ferrari · LiDAR Ride Height"):
    add_textbox(slide, 0.5, 7.0, 12, 0.4, text, 10, MUTED, False, PP_ALIGN.LEFT)
    add_textbox(slide, 11.5, 7.0, 1.5, 0.4, f"{slide.slide_id}", 10, MUTED, False, PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 1 — TITRE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg(slide)
add_textbox(slide, 0, 2.0, 13.333, 1.2, 'SCUDERIA FERRARI', 48, ROSSO, True, PP_ALIGN.CENTER, 'Calibri Light')
add_textbox(slide, 0, 3.2, 13.333, 0.8, 'LiDAR Ride Height', 32, WHITE, True, PP_ALIGN.CENTER, 'Calibri Light')
add_line(slide, 5.5, 4.2, 2.3, ROSSO)
add_textbox(slide, 0, 4.6, 13.333, 0.6, 'Cockpit Ingénieur · Télémétrie Temps Réel & IA', 18, MUTED, False, PP_ALIGN.CENTER)
add_textbox(slide, 0, 5.4, 13.333, 0.4, 'Projet Fefe — Démonstrateur Full-Stack', 14, MUTED, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — SOMMAIRE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.4, 12, 0.8, 'SOMMAIRE', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)
items = [
    '1. Vision & Contexte — Pourquoi la garde au sol est critique en F1',
    '2. Architecture Technique — Stack React / PHP / MySQL / IA',
    '3. Fonctionnalités Clés — Dashboard, Chatbot IA, Vocal',
    '4. Sécurité & Authentification — Google OAuth 2.0',
    '5. Design System — Thème Race Day, Dark/Light mode',
    '6. Démonstration Live',
    '7. Roadmap & Conclusion',
]
add_bullet_list(slide, 1.5, 1.8, 10, 5, items, 20, WHITE)
add_footer(slide)

# ============================================================
# SLIDE 3 — VISION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'VISION & CONTEXTE', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

add_textbox(slide, 0.8, 1.5, 5.5, 0.5, 'PROBLÉMATIQUE', 16, ROSSO, True)
items_left = [
    'La garde au sol (ride height) est critique en F1',
    'Un décrochage de plancher = perte d\'appui soudaine',
    'Les ingénieurs ont besoin de données en temps réel',
    'La complexité requiert une IA d\'assistance',
]
add_bullet_list(slide, 0.8, 2.2, 5.5, 3, items_left, 16, WHITE)

add_textbox(slide, 7.0, 1.5, 5.5, 0.5, 'NOTRE RÉPONSE', 16, ROSSO, True)
items_right = [
    'Dashboard temps réel avec LiDAR 100 Hz',
    'Résolution 0,1 mm sur la garde au sol',
    'Assistant IA Groq + Mistral intégré',
    'Contrôle vocal mains-libres',
    'Google OAuth sécurisé',
]
add_bullet_list(slide, 7.0, 2.2, 5.5, 3, items_right, 16, WHITE)

add_textbox(slide, 0.8, 5.8, 11.5, 0.8, '« La garde au sol se joue au millimètre. Notre LiDAR la mesure à 100 Hz. »', 16, MUTED, False, PP_ALIGN.CENTER, 'Calibri Light')
add_footer(slide)

# ============================================================
# SLIDE 4 — ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'ARCHITECTURE TECHNIQUE', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

layers = [
    ('React 19 · TanStack Router · Tailwind CSS 4 · TypeScript', 'FRONTEND'),
    ('Vite 7 · Hot Module Replacement', 'BUNDLER'),
    ('PHP 8.5 (Groq/Mistral) · Node.js Serverless (Vercel)', 'API'),
    ('MySQL · AlwaysData · Mesures · FAQ · LiDAR', 'BASE DE DONNÉES'),
    ('Groq (llama-3.3-70B) · Mistral (small-latest)', 'INTELLIGENCE ARTIFICIELLE'),
]
for i, (desc, label) in enumerate(layers):
    y = 1.6 + i * 1.05
    # Label
    add_textbox(slide, 1.0, y, 3, 0.4, label, 12, ROSSO, True)
    # Box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(y + 0.4), Inches(11), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = ROSSO
    shape.line.width = Pt(1)
    # Text in box
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

add_footer(slide)

# ============================================================
# SLIDE 5 — FLUX DE DONNÉES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'FLUX DE DONNÉES', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

flow = [
    '1. Capteurs IoT → Mesures (température, humidité, luminosité)',
    '2. LiDAR G2D → Distance (mm), réflectivité (%), luminosité (lux)',
    '3. MySQL → Stockage temps réel (AlwaysData)',
    '4. PHP API → Exposition REST des données',
    '5. React Dashboard → Affichage live (200 ms refresh)',
    '6. IA Groq/Mistral → Analyse contextuelle + conseils',
]
add_bullet_list(slide, 1.0, 1.6, 11, 3.5, flow, 18, WHITE)

add_textbox(slide, 1.0, 5.0, 11, 0.4, 'POINTS CLÉS', 16, ROSSO, True)
key_points = [
    '3 niveaux de fallback IA : Groq → Mistral → FAQ BDD → Local',
    'Proxy Vite PHP/Node pour le développement local',
    'Déploiement Vercel pour la production',
]
add_bullet_list(slide, 1.0, 5.5, 11, 1.5, key_points, 14, MUTED)
add_footer(slide)

# ============================================================
# SLIDE 6 — DASHBOARD
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'DASHBOARD COCKPIT INGÉNIEUR', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

add_textbox(slide, 0.8, 1.5, 5.5, 0.5, 'TÉLÉMÉTRIE LIVE', 16, ROSSO, True)
left = [
    'KPI : Downforce, Drag, Rake, Balance',
    'LiDAR AV/AR en mm (vue plongeante)',
    'Sliders de réglage setup',
    'Modèle aérodynamique temps réel',
    'Rafraîchissement 200 ms',
]
add_bullet_list(slide, 0.8, 2.1, 5.5, 3, left, 16, WHITE)

add_textbox(slide, 7.0, 1.5, 5.5, 0.5, 'VUES DISPONIBLES', 16, ROSSO, True)
right = [
    'Overview — KPIs + statut',
    'Architecture — Schéma système',
    'Telemetry — Setup + aéro 3D',
    'Simulators — Pit stop + stratégie',
    'IoT Devices — Capteurs G2D, G2E, LEDs',
]
add_bullet_list(slide, 7.0, 2.1, 5.5, 3, right, 16, WHITE)

add_textbox(slide, 0.8, 5.8, 11.5, 0.8, 'Statut piste : OPTIMAL · SUBOPTIMAL · CRITICAL', 18, GREEN, True, PP_ALIGN.CENTER)
add_footer(slide)

# ============================================================
# SLIDE 7 — IA CHATBOT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'ASSISTANT IA FERRARI', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

add_textbox(slide, 0.8, 1.5, 5.5, 0.5, 'CHATBOT INTELLIGENT', 16, ROSSO, True)
left = [
    'Contexte live (DB temps réel)',
    'Personnage : Ingénieur de Piste Ferrari',
    'Réponses en français, style italien 🇮🇹',
    'Rendu Markdown (gras, listes, code)',
    '3 niveaux de fallback',
]
add_bullet_list(slide, 0.8, 2.1, 5.5, 3, left, 16, WHITE)

add_textbox(slide, 7.0, 1.5, 5.5, 0.5, 'VOCAL & AUTH', 16, ROSSO, True)
right = [
    'Reconnaissance vocale (Web Speech API)',
    'Mode mains-libres : mot-clé « Ferrari »',
    'Synthèse vocale des réponses',
    'Google OAuth (connexion sécurisée)',
    'Fallback FERRARI/16',
]
add_bullet_list(slide, 7.0, 2.1, 5.5, 3, right, 16, WHITE)

# Example box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = PANEL
shape.line.color.rgb = ROSSO
shape.line.width = Pt(1)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '« Ciao ! Forza Ferrari ! La garde au sol est de 22,5 mm à l\'avant et 38,2 mm à l\'arrière. Le rake est dans la plage optimale. Le LiDAR confirme une bonne stabilité du plancher. »'
p.font.size = Pt(13)
p.font.color.rgb = GREEN
p.font.name = 'Calibri Light'
p.font.italic = True
p.alignment = PP_ALIGN.LEFT
add_footer(slide)

# ============================================================
# SLIDE 8 — MODE VOCAL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'MODE VOCAL MAINS-LIBRES', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

steps = [
    '1. Activation : bouton 🙌 dans le chatbot',
    '2. Écoute continue : reconnaissance vocale en boucle',
    '3. Mot-clé : « Ferrari » (ou « Hey Ferrari », « OK Ferrari »)',
    '4. Bip sonore (Web Audio API) : confirmation de détection',
    '5. Question : dictée après le bip',
    '6. Réponse : l\'IA répond + synthèse vocale',
]
add_bullet_list(slide, 1.0, 1.6, 11, 3.5, steps, 20, WHITE)

add_textbox(slide, 1.0, 5.2, 11, 0.4, 'SCÉNARIOS', 16, ROSSO, True)
scenarios = [
    '« Ferrari quelle est la garde au sol ? » → réponse directe immédiate',
    '« Ferrari » [pause] « analyse le setup » → deux temps, attente active',
]
add_bullet_list(slide, 1.0, 5.7, 11, 1, scenarios, 14, MUTED)
add_footer(slide)

# ============================================================
# SLIDE 9 — API IA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'API IA — 3 NIVEAUX DE FALLBACK', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

fallbacks = [
    ('GROQ', 'llama-3.3-70b-versatile · Latence < 500 ms', ROSSO),
    ('MISTRAL (fallback)', 'mistral-small-latest · Deep reasoning', AMBER),
    ('FAQ BDD (fallback)', 'Recherche MySQL · Questions approuvées', AMBER),
    ('MODE HORS-LIGNE', 'Réponses codées · Données simulées', GREEN),
]
for i, (name, desc, color) in enumerate(fallbacks):
    y = 1.6 + i * 1.3
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(y), Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = 'Calibri'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = MUTED
    p2.font.name = 'Calibri'

    if i < 3:
        # Arrow down
        add_textbox(slide, 6.2, y + 0.9, 1, 0.4, '↓', 20, ROSSO, True, PP_ALIGN.CENTER)

add_footer(slide)

# ============================================================
# SLIDE 10 — SÉCURITÉ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'GOOGLE OAUTH 2.0 & SÉCURITÉ', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

add_textbox(slide, 0.8, 1.5, 5.5, 0.5, 'FLUX D\'AUTHENTIFICATION', 16, ROSSO, True)
left = [
    '1. Clic sur « Google Sign-In »',
    '2. Google vérifie l\'identité',
    '3. JWT décodé côté client',
    '4. Session persistée (localStorage)',
    '5. Dashboard protégé par useAuth()',
]
add_bullet_list(slide, 0.8, 2.1, 5.5, 3, left, 16, WHITE)

add_textbox(slide, 7.0, 1.5, 5.5, 0.5, 'SÉCURITÉ', 16, ROSSO, True)
right = [
    'Clés API jamais exposées au client',
    '.env.local exclu du Git (.gitignore)',
    'Rate limiting : 10 req/min par IP',
    'Fallback local FERRARI/16',
    'CORS headers sur l\'API',
]
add_bullet_list(slide, 7.0, 2.1, 5.5, 3, right, 16, WHITE)

# Alert box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xDC, 0x00, 0x00)
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '🔒  Les clés Groq/Mistral sont stockées côté serveur (PHP) et JAMAIS envoyées au client. Le frontend ne voit que les réponses.'
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.bold = True
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER
add_footer(slide)

# ============================================================
# SLIDE 11 — DESIGN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'DESIGN SYSTEM — THÈME « RACE DAY »', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

# Dark mode box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK
shape.line.color.rgb = ROSSO
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '🌙  MODE SOMBRE (défaut)'
p.font.size = Pt(18)
p.font.color.rgb = ROSSO
p.font.bold = True
p.font.name = 'Calibri'
colors = ['Fond : #0A0A0A', 'Panels : #121212', 'Texte : #F5F5F5', 'Accent : Rosso Corsa #DC0000', 'Sensor : Vert #00FF41']
for c in colors:
    p2 = tf.add_paragraph()
    p2.text = '   ' + c
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.font.name = 'Calibri'
    p2.space_after = Pt(4)

# Light mode box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.6), Inches(5.5), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
shape.line.color.rgb = ROSSO
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '☀️  MODE CLAIR (toggle)'
p.font.size = Pt(18)
p.font.color.rgb = ROSSO
p.font.bold = True
p.font.name = 'Calibri'
colors = ['Fond : #F5F5F5', 'Panels : #FFFFFF', 'Texte : #1A1A1A', 'Accent : Rosso Corsa #DC0000', 'Sensor : Vert #0A7A2E']
for c in colors:
    p2 = tf.add_paragraph()
    p2.text = '   ' + c
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    p2.font.name = 'Calibri'
    p2.space_after = Pt(4)

add_textbox(slide, 0.8, 5.8, 11.5, 0.6, 'Tailwind CSS 4 · JetBrains Mono · Inter · Animations CSS personnalisées', 14, MUTED, False, PP_ALIGN.CENTER)
add_footer(slide)

# ============================================================
# SLIDE 12 — DÉMO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'DÉMONSTRATION LIVE', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

steps = [
    '1. Vitrine publique — Storytelling Ferrari, circuit Monza interactif au scroll',
    '2. Connexion Google — OAuth ou identifiants FERRARI / 16',
    '3. Cockpit Dashboard — Télémétrie live, sliders setup, KPI',
    '4. Assistant IA — Posez une question technique, réponse Groq/Mistral',
    '5. Mode vocal — Dites « Ferrari » + votre question, mains-libres',
    '6. Toggle ☀️/🌙 — Mode clair/sombre en un clic',
]
add_bullet_list(slide, 1.5, 1.6, 10, 4.5, steps, 22, WHITE)

add_textbox(slide, 0.8, 6.4, 11.5, 0.6, 'npm run dev:php', 24, GREEN, True, PP_ALIGN.CENTER, 'Consolas')
add_footer(slide)

# ============================================================
# SLIDE 13 — ROADMAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'ROADMAP & PERSPECTIVES', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

add_textbox(slide, 0.8, 1.5, 5.5, 0.5, 'COURT TERME', 16, ROSSO, True)
left = [
    'Streaming des réponses IA (token par token)',
    'Historique des conversations (BDD)',
    'Quick Actions / prompts prédéfinis',
    'Notifications push (Safety Car, drapeaux)',
]
add_bullet_list(slide, 0.8, 2.1, 5.5, 3, left, 16, WHITE)

add_textbox(slide, 7.0, 1.5, 5.5, 0.5, 'MOYEN TERME', 16, ROSSO, True)
right = [
    'Données LiDAR réelles (API F1)',
    'Mode multi-écrans (pit wall)',
    'Analyse prédictive (ML)',
    'Export PDF des rapports setup',
    'PWA hors-ligne',
]
add_bullet_list(slide, 7.0, 2.1, 5.5, 3, right, 16, WHITE)

add_textbox(slide, 0.8, 5.8, 11.5, 0.6, '« Un véritable outil d\'ingénierie piste, pas juste un dashboard. »', 16, MUTED, False, PP_ALIGN.CENTER, 'Calibri Light')
add_footer(slide)

# ============================================================
# SLIDE 14 — CONCLUSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0, 0.8, 13.333, 1.2, 'FORZA FERRARI!', 52, ROSSO, True, PP_ALIGN.CENTER, 'Calibri Light')

items = [
    'Télémétrie temps réel — LiDAR, aéro, setup',
    'IA Groq + Mistral — Assistant contextuel 3 niveaux',
    'Vocal mains-libres — Mot-clé « Ferrari », Web Speech API',
    'Google OAuth — Connexion sécurisée, localStorage',
    'Design Ferrari — Dark/Light mode, toggle navbar',
    'Full-stack — React 19, PHP 8.5, MySQL, Vercel',
]
add_bullet_list(slide, 2.5, 2.4, 8, 3.5, items, 20, WHITE)

add_line(slide, 5.5, 6.0, 2.3, ROSSO)
add_textbox(slide, 0, 6.2, 13.333, 0.6, 'github.com/evanmse/Fefe_project', 16, WHITE, False, PP_ALIGN.CENTER, 'Consolas')
add_textbox(slide, 0, 6.7, 13.333, 0.4, 'Questions ?', 14, MUTED, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 15 — ANNEXE STACK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.8, 'ANNEXE — STACK DÉTAILLÉE', 28, ROSSO, True, PP_ALIGN.LEFT, 'Calibri Light')
add_red_bar(slide)

for col_idx, (title, techs) in enumerate([
    ('FRONTEND', ['React 19', 'TanStack Router', 'Tailwind CSS 4', 'TypeScript 5.7', 'Web Speech API', 'Web Audio API']),
    ('BACKEND', ['PHP 8.5 (API REST)', 'Node.js (Vercel)', 'MySQL (AlwaysData)', 'cURL (Groq/Mistral)', 'PDO (DB)']),
    ('DEVOPS', ['Vite 7', 'Vercel', 'GitHub (2 remotes)', 'gh CLI', 'npm scripts']),
]):
    x = 0.8 + col_idx * 4.2
    add_textbox(slide, x, 1.6, 3.5, 0.5, title, 16, ROSSO, True)
    add_bullet_list(slide, x, 2.2, 3.5, 3, techs, 14, WHITE)

add_footer(slide)

# ============================================================
# SAVE
# ============================================================
output = '/Users/s.sy/Fefe_project/Fefe_project/presentation.pptx'
prs.save(output)
print(f'✅ PowerPoint saved: {output}')
print(f'   Slides: {len(prs.slides)}')
